#!/usr/bin/env python3
"""Genera una plantilla PowerPoint de ejemplo con estilos corporativos."""

import sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).parent.parent / "src"))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

AZUL = RGBColor(0x00, 0x3F, 0x87)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
GRIS = RGBColor(0xF5, 0xF5, 0xF5)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Slide 1: Portada ──────────────────────────────────────────────────────────
layout = prs.slide_layouts[0]  # Title Slide
slide = prs.slides.add_slide(layout)

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = AZUL

title = slide.shapes.title
title.text = "Título de la Presentación"
title.text_frame.paragraphs[0].font.color.rgb = BLANCO
title.text_frame.paragraphs[0].font.size = Pt(40)
title.text_frame.paragraphs[0].font.bold = True

subtitle = slide.placeholders[1]
subtitle.text = "Subtítulo o fecha"
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
subtitle.text_frame.paragraphs[0].font.size = Pt(20)

# ── Slide 2: Contenido (ejemplo de layout de bullets) ─────────────────────────
layout2 = prs.slide_layouts[1]  # Title and Content
slide2 = prs.slides.add_slide(layout2)

slide2.shapes.title.text = "Ejemplo de Diapositiva de Contenido"
slide2.shapes.title.text_frame.paragraphs[0].font.color.rgb = AZUL
slide2.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
slide2.shapes.title.text_frame.paragraphs[0].font.bold = True

content_ph = slide2.placeholders[1]
tf = content_ph.text_frame
tf.text = "Punto principal 1"
p2 = tf.add_paragraph()
p2.text = "Punto principal 2"
p3 = tf.add_paragraph()
p3.text = "Punto principal 3"

out_path = Path(__file__).parent / "plantilla_ejemplo.pptx"
prs.save(out_path)
print(f"Plantilla creada: {out_path}")
