"""Builds PowerPoint presentations from a template and LLM-generated slide plan."""

import copy
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def _find_layout(prs: Presentation, layout_name: str):
    """Returns the layout matching the name, or the first layout as fallback."""
    for layout in prs.slide_layouts:
        if layout.name.lower() == layout_name.lower():
            return layout
    return prs.slide_layouts[0]


def _set_text(placeholder, text: str, bold: bool = False, font_size: int = None):
    tf = placeholder.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(text)
    if bold:
        run.font.bold = True
    if font_size:
        run.font.size = Pt(font_size)


def build_presentation(template_path: str, slide_plan: list[dict], output_path: str) -> str:
    """
    Builds a PowerPoint from a template using a slide plan.

    slide_plan is a list of dicts, each with:
        {
            "layout": "Title Slide",          # layout name from template
            "title": "Slide title",
            "subtitle": "Optional subtitle",
            "content": ["bullet 1", "bullet 2"],  # optional
            "notes": "Speaker notes"          # optional
        }
    """
    prs = Presentation(template_path)

    # Remove existing slides (keep template master/layouts, clear slides)
    xml_slides = prs.slides._sldIdLst
    for slide in list(prs.slides):
        rId = prs.slides._sldIdLst[0].get("r:id")
        # We rebuild from scratch using add_slide
        break

    # Clear all existing slides
    slide_ids = list(prs.slides._sldIdLst)
    for s_id in slide_ids:
        rId = s_id.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        prs.part.drop_rel(rId)
        xml_slides.remove(s_id)

    for slide_data in slide_plan:
        layout_name = slide_data.get("layout", "")
        layout = _find_layout(prs, layout_name)
        slide = prs.slides.add_slide(layout)

        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            name = ph.name.lower()

            if idx == 0 or "title" in name:
                title_text = slide_data.get("title", "")
                if title_text:
                    _set_text(ph, title_text, bold=True)

            elif idx == 1 and "subtitle" in name:
                subtitle = slide_data.get("subtitle", "")
                if subtitle:
                    _set_text(ph, subtitle)

            elif idx == 1 or "content" in name or "body" in name:
                content = slide_data.get("content", [])
                if content:
                    tf = ph.text_frame
                    tf.clear()
                    for i, bullet in enumerate(content):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = str(bullet)
                        p.level = 0

        # Add speaker notes
        notes_text = slide_data.get("notes", "")
        if notes_text:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.text = notes_text

    out_path = Path(output_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return str(out_path)
