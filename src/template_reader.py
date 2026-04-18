"""Reads PowerPoint templates and extracts structure information."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
import json


def read_template(template_path: str) -> dict:
    """
    Reads a PowerPoint template and returns its structure:
    slide layouts, placeholders, fonts, colors, and dimensions.
    """
    path = Path(template_path)
    if not path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    prs = Presentation(template_path)

    width_in = prs.slide_width.inches
    height_in = prs.slide_height.inches

    layouts = []
    for layout in prs.slide_layouts:
        placeholders = []
        for ph in layout.placeholders:
            placeholders.append({
                "idx": ph.placeholder_format.idx,
                "name": ph.name,
                "type": str(ph.placeholder_format.type),
            })
        layouts.append({
            "name": layout.name,
            "placeholders": placeholders,
        })

    existing_slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        existing_slides.append({
            "slide_number": i + 1,
            "layout": slide.slide_layout.name,
            "texts": texts,
        })

    return {
        "dimensions": {"width_inches": width_in, "height_inches": height_in},
        "layouts": layouts,
        "existing_slides": existing_slides,
        "total_slides": len(prs.slides),
    }


def summarize_template(template_info: dict) -> str:
    """Converts template info into a readable text for the LLM."""
    lines = [
        f"Dimensiones: {template_info['dimensions']['width_inches']:.1f}\" x {template_info['dimensions']['height_inches']:.1f}\"",
        f"Diapositivas existentes: {template_info['total_slides']}",
        "",
        "Layouts disponibles:",
    ]
    for layout in template_info["layouts"]:
        ph_names = ", ".join(ph["name"] for ph in layout["placeholders"]) or "ninguno"
        lines.append(f"  - {layout['name']} | Placeholders: {ph_names}")

    if template_info["existing_slides"]:
        lines.append("")
        lines.append("Diapositivas en la plantilla:")
        for slide in template_info["existing_slides"]:
            lines.append(f"  Slide {slide['slide_number']} ({slide['layout']}): {' / '.join(slide['texts'][:3])}")

    return "\n".join(lines)
