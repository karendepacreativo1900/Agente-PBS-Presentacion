"""Tests for presentation_builder module."""

import sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).parent.parent / "src"))

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from presentation_builder import build_presentation


@pytest.fixture
def template_path(tmp_path):
    """Creates a minimal template .pptx for testing."""
    prs = Presentation()
    out = tmp_path / "template.pptx"
    prs.save(str(out))
    return str(out)


def test_build_creates_file(template_path, tmp_path):
    plan = [
        {"layout": "Title Slide", "title": "Mi Presentación", "subtitle": "2024"},
        {"layout": "Title and Content", "title": "Datos", "content": ["Punto A", "Punto B"]},
    ]
    output = str(tmp_path / "output.pptx")
    result = build_presentation(template_path, plan, output)
    assert Path(result).exists()


def test_build_correct_slide_count(template_path, tmp_path):
    plan = [
        {"layout": "Title Slide", "title": "Portada"},
        {"layout": "Title and Content", "title": "Contenido 1", "content": ["A"]},
        {"layout": "Title and Content", "title": "Contenido 2", "content": ["B"]},
    ]
    output = str(tmp_path / "output2.pptx")
    build_presentation(template_path, plan, output)
    prs = Presentation(output)
    assert len(prs.slides) == 3
